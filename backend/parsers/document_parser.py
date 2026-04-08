"""
parsers/document_parser.py  —  Universal Document Parser
=========================================================
Supports: PDF · Excel/XLSX/XLS · CSV · JSON · XML · DOCX · PPTX
          HTML · Circulars · Notices · Press Releases · Reports

Handles both URLs (downloaded) and uploaded byte content.
All output is normalized to a common schema.
"""

import io, json, re, csv, xml.etree.ElementTree as ET
from typing import Optional, Union
from pathlib import Path

# ── Output schema ─────────────────────────────────────────────────────────────
def _doc(source: str, doc_type: str, title: str, text: str,
         tables: list = None, metadata: dict = None, raw_pages: list = None) -> dict:
    return {
        "source":     source,
        "type":       doc_type,
        "title":      title,
        "text":       text[:50000],   # cap at 50k chars
        "word_count": len(text.split()),
        "tables":     tables or [],
        "metadata":   metadata or {},
        "raw_pages":  raw_pages or [],
        "status":     "parsed",
    }


# ── PDF Parser ────────────────────────────────────────────────────────────────

def parse_pdf_bytes(data: bytes, source: str = "upload") -> dict:
    """Parse PDF using PyMuPDF (fitz). Extracts text + tables from all pages."""
    try:
        import fitz   # PyMuPDF
        doc    = fitz.open(stream=data, filetype="pdf")
        pages  = []
        tables = []
        for i, page in enumerate(doc):
            text = page.get_text("text")
            pages.append({"page": i+1, "text": text.strip()})
            # Extract tables via find_tables
            try:
                tab = page.find_tables()
                for t in tab.tables:
                    rows = t.extract()
                    if rows:
                        tables.append({"page": i+1, "rows": rows})
            except Exception:
                pass
        full_text = "\n\n".join(p["text"] for p in pages)
        meta  = doc.metadata or {}
        title = meta.get("title") or _guess_title(full_text)
        doc.close()
        return _doc(source, "PDF", title, full_text,
                    tables=tables, metadata=meta, raw_pages=pages)
    except ImportError:
        return _pdf_fallback(data, source)
    except Exception as ex:
        return {"error": str(ex), "source": source, "type": "PDF", "status": "failed"}


def _pdf_fallback(data: bytes, source: str) -> dict:
    """Fallback: extract text from PDF without PyMuPDF using basic pattern matching."""
    try:
        text = data.decode("latin-1", errors="replace")
        # Extract readable strings (printable ASCII sequences > 4 chars)
        chunks = re.findall(r'[!-~\s]{5,}', text)
        readable = " ".join(c.strip() for c in chunks if len(c.strip()) > 4)
        # Clean up
        readable = re.sub(r'\s+', ' ', readable).strip()
        title = _guess_title(readable)
        return _doc(source, "PDF", title, readable,
                    metadata={"note": "PyMuPDF not installed — basic extraction used"})
    except Exception as ex:
        return {"error": str(ex), "source": source, "type": "PDF", "status": "failed"}


# ── Excel / XLSX / XLS Parser ────────────────────────────────────────────────

def parse_excel_bytes(data: bytes, source: str = "upload",
                      filename: str = "file.xlsx") -> dict:
    """Parse Excel files. Returns all sheets as tables + combined text."""
    try:
        import openpyxl
        wb     = openpyxl.load_workbook(io.BytesIO(data), data_only=True)
        tables = []
        text_parts = []
        for sheet_name in wb.sheetnames:
            ws   = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                clean = [str(c) if c is not None else "" for c in row]
                if any(c.strip() for c in clean):
                    rows.append(clean)
            if rows:
                tables.append({"sheet": sheet_name, "rows": rows})
                # Convert to text
                for row in rows:
                    text_parts.append("  |  ".join(row))
        full_text = "\n".join(text_parts)
        title = _guess_title(full_text) or filename
        return _doc(source, "Excel", title, full_text,
                    tables=tables, metadata={"sheets": wb.sheetnames, "filename": filename})
    except ImportError:
        pass
    # Fallback to xlrd for .xls
    try:
        import xlrd
        book   = xlrd.open_workbook(file_contents=data)
        tables = []; text_parts = []
        for sheet in book.sheets():
            rows = [
                [str(sheet.cell_value(r, c)) for c in range(sheet.ncols)]
                for r in range(sheet.nrows)
            ]
            tables.append({"sheet": sheet.name, "rows": rows})
            for row in rows:
                text_parts.append("  |  ".join(row))
        full_text = "\n".join(text_parts)
        return _doc(source, "Excel", filename, full_text, tables=tables)
    except Exception as ex:
        return {"error": str(ex), "source": source, "type": "Excel", "status": "failed"}


# ── CSV Parser ────────────────────────────────────────────────────────────────

def parse_csv_bytes(data: bytes, source: str = "upload") -> dict:
    text_io  = io.StringIO(data.decode("utf-8", errors="replace"))
    reader   = csv.reader(text_io)
    rows     = list(reader)
    headers  = rows[0] if rows else []
    records  = rows[1:] if len(rows) > 1 else []
    # Build text representation
    text = "\n".join("  |  ".join(r) for r in rows[:200])
    return _doc(source, "CSV", f"CSV Dataset ({len(records)} rows)", text,
                tables=[{"rows": rows[:500]}],
                metadata={"row_count": len(records), "columns": headers,
                          "column_count": len(headers)})


# ── JSON Parser ───────────────────────────────────────────────────────────────

def parse_json_bytes(data: bytes, source: str = "upload") -> dict:
    try:
        obj      = json.loads(data.decode("utf-8", errors="replace"))
        text     = json.dumps(obj, indent=2)[:20000]
        is_list  = isinstance(obj, list)
        meta = {
            "is_array":     is_list,
            "record_count": len(obj) if is_list else 1,
            "keys":         list(obj[0].keys()) if is_list and obj and isinstance(obj[0], dict)
                            else list(obj.keys()) if isinstance(obj, dict) else [],
        }
        title = _guess_title(text) or "JSON Dataset"
        return _doc(source, "JSON", title, text, metadata=meta)
    except Exception as ex:
        return {"error": str(ex), "source": source, "type": "JSON", "status": "failed"}


# ── XML Parser ────────────────────────────────────────────────────────────────

def parse_xml_bytes(data: bytes, source: str = "upload") -> dict:
    try:
        root  = ET.fromstring(data)
        items = []
        def _walk(el, depth=0):
            text = (el.text or "").strip()
            if text:
                items.append(f"{'  '*depth}<{el.tag}> {text}")
            for child in el:
                _walk(child, depth+1)
        _walk(root)
        full_text = "\n".join(items)
        return _doc(source, "XML", root.tag, full_text,
                    metadata={"root_tag": root.tag, "attributes": root.attrib})
    except Exception as ex:
        return {"error": str(ex), "source": source, "type": "XML", "status": "failed"}


# ── DOCX Parser ───────────────────────────────────────────────────────────────

def parse_docx_bytes(data: bytes, source: str = "upload") -> dict:
    try:
        from docx import Document
        doc   = Document(io.BytesIO(data))
        paras = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        tables = []
        for t in doc.tables:
            rows = [[cell.text for cell in row.cells] for row in t.rows]
            tables.append({"rows": rows})
        full_text = "\n\n".join(paras)
        title = paras[0][:100] if paras else "Word Document"
        return _doc(source, "DOCX", title, full_text, tables=tables,
                    metadata={"paragraph_count": len(paras), "table_count": len(tables)})
    except ImportError:
        return {"error": "python-docx not installed", "source": source,
                "type": "DOCX", "status": "failed"}
    except Exception as ex:
        return {"error": str(ex), "source": source, "type": "DOCX", "status": "failed"}


# ── PPTX Parser ───────────────────────────────────────────────────────────────

def parse_pptx_bytes(data: bytes, source: str = "upload") -> dict:
    try:
        from pptx import Presentation
        prs   = Presentation(io.BytesIO(data))
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            slides.append({"slide": i, "content": "\n".join(texts)})
        full_text = "\n\n".join(f"--- Slide {s['slide']} ---\n{s['content']}" for s in slides)
        return _doc(source, "PPTX", f"Presentation ({len(slides)} slides)", full_text,
                    raw_pages=slides, metadata={"slide_count": len(slides)})
    except ImportError:
        return {"error": "python-pptx not installed", "source": source,
                "type": "PPTX", "status": "failed"}
    except Exception as ex:
        return {"error": str(ex), "source": source, "type": "PPTX", "status": "failed"}


# ── HTML Parser ───────────────────────────────────────────────────────────────

def parse_html(html: str, source: str = "url") -> dict:
    from bs4 import BeautifulSoup
    soup   = BeautifulSoup(html, "lxml")
    # Remove scripts/styles
    for tag in soup(["script", "style", "nav", "footer", "header"]):
        tag.decompose()
    title_el = soup.find("title")
    title    = title_el.get_text(strip=True) if title_el else ""
    # Extract all text
    text = soup.get_text(" ", strip=True)
    text = re.sub(r'\s+', ' ', text).strip()
    # Extract tables
    tables = []
    for tbl in soup.find_all("table"):
        rows = []
        for tr in tbl.find_all("tr"):
            cells = [td.get_text(strip=True) for td in tr.find_all(["td","th"])]
            if cells: rows.append(cells)
        if rows: tables.append({"rows": rows})
    return _doc(source, "HTML", title, text, tables=tables)


# ── Circular / Notice / Press Release detector ───────────────────────────────

CIRCULAR_KEYWORDS   = ["circular", "notification", "office memorandum", "o.m.", "corrigendum"]
NOTICE_KEYWORDS     = ["notice", "public notice", "tender notice", "vacancy notice"]
PRESS_REL_KEYWORDS  = ["press release", "press note", "press information bureau",
                       "ministry of", "government of india announces"]
REPORT_KEYWORDS     = ["annual report", "report", "survey", "committee report", "white paper"]

def classify_document(text: str, title: str = "") -> str:
    """Classify document type from its content."""
    combined = (title + " " + text[:1000]).lower()
    if any(k in combined for k in CIRCULAR_KEYWORDS):   return "Circular"
    if any(k in combined for k in NOTICE_KEYWORDS):     return "Notice"
    if any(k in combined for k in PRESS_REL_KEYWORDS):  return "Press Release"
    if any(k in combined for k in REPORT_KEYWORDS):     return "Report"
    return "Document"


# ── URL-based dispatcher ──────────────────────────────────────────────────────

async def parse_from_url(url: str) -> dict:
    """Download and parse any supported document from a URL."""
    from utils.session import fetch_binary, fetch
    url_lower = url.lower()
    try:
        if any(url_lower.endswith(ext) for ext in [".pdf"]):
            data  = await fetch_binary(url)
            result = parse_pdf_bytes(data, source=url)
        elif any(url_lower.endswith(ext) for ext in [".xlsx", ".xls"]):
            data  = await fetch_binary(url)
            filename = url.split("/")[-1]
            result = parse_excel_bytes(data, source=url, filename=filename)
        elif url_lower.endswith(".csv"):
            data  = await fetch_binary(url)
            result = parse_csv_bytes(data, source=url)
        elif url_lower.endswith(".json"):
            data  = await fetch_binary(url)
            result = parse_json_bytes(data, source=url)
        elif url_lower.endswith(".xml"):
            data  = await fetch_binary(url)
            result = parse_xml_bytes(data, source=url)
        elif any(url_lower.endswith(ext) for ext in [".docx", ".doc"]):
            data  = await fetch_binary(url)
            result = parse_docx_bytes(data, source=url)
        elif url_lower.endswith(".pptx"):
            data  = await fetch_binary(url)
            result = parse_pptx_bytes(data, source=url)
        else:
            # Treat as HTML
            r = await fetch(url)
            result = parse_html(r.text, source=url)

        # Classify doc type
        if result.get("status") == "parsed":
            result["doc_class"] = classify_document(
                result.get("text",""), result.get("title","")
            )
        return result
    except Exception as ex:
        return {"error": str(ex), "url": url, "status": "failed"}


def parse_bytes(data: bytes, filename: str, source: str = "upload") -> dict:
    """Auto-detect format from filename and parse bytes."""
    fn = filename.lower()
    if fn.endswith(".pdf"):   return parse_pdf_bytes(data, source)
    if fn.endswith((".xlsx",".xls")): return parse_excel_bytes(data, source, filename)
    if fn.endswith(".csv"):   return parse_csv_bytes(data, source)
    if fn.endswith(".json"):  return parse_json_bytes(data, source)
    if fn.endswith(".xml"):   return parse_xml_bytes(data, source)
    if fn.endswith((".docx",".doc")): return parse_docx_bytes(data, source)
    if fn.endswith(".pptx"):  return parse_pptx_bytes(data, source)
    if fn.endswith((".html",".htm")):
        return parse_html(data.decode("utf-8", errors="replace"), source)
    # Unknown — try text
    try:
        text = data.decode("utf-8", errors="replace")
        return _doc(source, "Text", _guess_title(text), text)
    except Exception as ex:
        return {"error": f"Unsupported format: {filename}", "status": "failed"}


def _guess_title(text: str, max_len: int = 100) -> str:
    """Guess document title from first non-empty line."""
    for line in text.split("\n"):
        clean = line.strip()
        if len(clean) > 8:
            return clean[:max_len]
    return "Untitled Document"
